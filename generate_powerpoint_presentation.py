"""
generate_powerpoint_presentation.py – Generates an expanded 10-slide, highly professional,
visually appealing widescreen PowerPoint presentation (.pptx) embedding visual diagrams.
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Dark Enterprise Theme) ───────────────────────────
C_DARK_BG     = RGBColor(15, 23, 42)     # #0F172A Slate 900
C_CARD_BG     = RGBColor(30, 41, 59)     # #1E293B Slate 800
C_CARD_BORDER = RGBColor(71, 85, 105)  # #475569 Slate 600
C_PRIMARY     = RGBColor(248, 250, 252) # #F8FAFC White/Slate 50
C_ACCENT      = RGBColor(99, 102, 241)   # #6366F1 Indigo 500
C_SUCCESS     = RGBColor(16, 185, 129)   # #10B981 Emerald 500
C_WARNING     = RGBColor(245, 158, 11)   # #F59E0B Amber 500
C_MUTED       = RGBColor(148, 163, 184)  # #94A3B8 Slate 400


def apply_slide_background(slide):
    """Sets a dark gradient background on the slide using a full rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BG
    bg.line.fill.background()
    return bg


def add_header(slide, title_text, subtitle_text=None):
    """Adds a standardized slide header with title and subtitle."""
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = "Arial"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_MUTED
        p2.space_before = Pt(3)


def add_card(slide, left, top, width, height, title="", body="", bg_color=C_CARD_BG, border_color=C_CARD_BORDER, title_color=C_ACCENT):
    """Adds a rounded rectangle card with title and body text."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    if title:
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.name = "Arial"
        p0.font.size = Pt(13.5)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        p0.space_after = Pt(6)
        
        if body:
            p1 = tf.add_paragraph()
            p1.text = body
            p1.font.name = "Arial"
            p1.font.size = Pt(10)
            p1.font.color.rgb = C_PRIMARY
            p1.space_before = Pt(3)
    else:
        if body:
            p0 = tf.paragraphs[0]
            p0.text = body
            p0.font.name = "Arial"
            p0.font.size = Pt(10)
            p0.font.color.rgb = C_PRIMARY

    return card


def build_presentation(filename="ECommerce_Agentic_AI_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE SLIDE
    # ═════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.0), Inches(4.5))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "E-COMMERCE DATA ENGINEERING &\nAGENTIC AI PLATFORM"
    p.font.name = "Arial"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    p.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "End-to-End Medallion Data Engineering Pipeline + 9-Agent Autonomous AI Engine"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_MUTED
    p2.space_after = Pt(24)

    p3 = tf1.add_paragraph()
    p3.text = "DEVELOPER: Student Project  |  TECH STACK: PostgreSQL • Apache Airflow • Gemini AI • Streamlit • Docker"
    p3.font.name = "Arial"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = C_SUCCESS

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEM STATEMENT VS SOLUTION
    # ═════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide2)
    add_header(slide2, "1. Business Problem vs. My Solution", "Why this platform was built and what business value it delivers.")

    add_card(
        slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        title="❌ THE BUSINESS PROBLEM",
        body=(
            "1. Raw Data Messiness:\n"
            "   E-Commerce stores process thousands of daily orders. Raw data is scattered across files with missing values, invalid dates, and currency errors.\n\n"
            "2. Delay & Reliance on Engineers:\n"
            "   Whenever managers need answers (e.g. 'What were sales in California?'), they must wait days for data engineers to write custom database code.\n\n"
            "3. Lack of Actionable Insights:\n"
            "   Static spreadsheets show numbers, but fail to explain WHY sales dropped or WHAT management should do next."
        ),
        title_color=RGBColor(239, 68, 68)
    )

    add_card(
        slide2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        title="✅ MY INTEGRATED SOLUTION",
        body=(
            "1. Automated Medallion Pipeline (Layer 1):\n"
            "   A 3-tier pipeline (Bronze → Silver → Gold) powered by PostgreSQL and Airflow that automatically cleans data, tracks location history (SCD2), and builds Star Schema tables.\n\n"
            "2. Autonomous Agentic AI Layer (Layer 2):\n"
            "   A team of 9 specialized AI agents connected to the Gold database that turns natural English questions into instant SQL queries, Plotly charts, and business recommendations!\n\n"
            "3. Zero Dependency on Data Teams:\n"
            "   Managers get instant answers in seconds!"
        ),
        title_color=C_SUCCESS
    )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 3: MEDALLION PIPELINE ARCHITECTURE (WITH DIAGRAM)
    # ═════════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide3)
    add_header(slide3, "2. Data Engineering: Medallion Architecture", "Refining raw incoming data into pure, analytics-ready Star Schema tables.")

    # Embed visual diagram image if available
    if os.path.exists("medallion_diagram.png"):
        slide3.shapes.add_picture("medallion_diagram.png", Inches(0.8), Inches(1.5), width=Inches(11.733))

    # Text summary below diagram
    add_card(
        slide3, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.3),
        title="MEDALLION PIPELINE BREAKDOWN",
        body=(
            "• Bronze Layer (🟫 Raw): Append-only storage for original customer order feeds. Preserves raw state with batch IDs and timestamps.\n"
            "• Silver Layer (🥈 Cleansed): Standardizes date formats, cleans currencies, removes duplicates, translates categories to English, and tracks SCD Type 2.\n"
            "• Gold Layer (🥇 Analytics Star Schema): Optimized dimensional star schema (fact_sales, dim_customer, dim_product, dim_seller, dim_date, revenue_mart)."
        ),
        title_color=C_ACCENT
    )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 4: SCD TYPE 1/2 & DATA QUALITY
    # ═════════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide4)
    add_header(slide4, "3. Historical Tracking (SCD) & Data Quality Rules", "Ensuring point-in-time revenue accuracy and enforcing data validation.")

    add_card(
        slide4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        title="🔄 SLOWLY CHANGING DIMENSIONS (SCD)",
        body=(
            "SCD Type 1 (Overwrite):\n"
            "• Target: gold.dim_product\n"
            "• Used for basic product updates (weight, size category) where history isn't needed.\n"
            "• SQL: ON CONFLICT (product_id) DO UPDATE SET...\n\n"
            "SCD Type 2 (Historical Versioning):\n"
            "• Target: gold.dim_customer & gold.dim_seller\n"
            "• Tracks relocations across cities/states.\n"
            "• Mechanism: Expire old record (is_current = FALSE, end_date = NOW()) and INSERT new record version (is_current = TRUE).\n"
            "• Guarantees past sales stay credited to original regions!"
        ),
        title_color=C_ACCENT
    )

    add_card(
        slide4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        title="🛡️ DATA QUALITY & METADATA AUDITING",
        body=(
            "1. Primary Key Non-Null Validation:\n"
            "   Enforces non-null rules on customer_id, order_id, and product_id before processing.\n\n"
            "2. Range & Value Constraints:\n"
            "   Prices, freight fees, and payment values must be strictly > $0.00.\n\n"
            "3. Referential Integrity Check:\n"
            "   Fact sales records must resolve strictly to active dimension surrogate keys.\n\n"
            "4. Metadata Audit Log:\n"
            "   Tracks all runs in metadata.pipeline_runs, failed records in metadata.dq_results, and watermarks."
        ),
        title_color=C_SUCCESS
    )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 5: DATA SCIENCE & ML ENGINE (WITH CHART)
    # ═════════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide5)
    add_header(slide5, "4. Data Science & ML Engine", "Regression trends, moving averages, Z-score anomaly detection, and linear forecasting.")

    # Left: Explanation card
    add_card(
        slide5, Inches(0.8), Inches(1.6), Inches(5.4), Inches(5.2),
        title="🧠 STATISTICAL & ML CAPABILITIES",
        body=(
            "1. Linear Trend Regression:\n"
            "   Computes regression slopes on monthly revenue time-series to determine whether sales trajectories are increasing or decreasing.\n\n"
            "2. Z-Score Anomaly Detection:\n"
            "   Calculates statistical Z-scores across monthly revenue. Months where |Z| > 2.0 are flagged as revenue anomalies.\n\n"
            "3. Revenue Forecasting:\n"
            "   Extrapolates sales for 1 to 12 months ahead using linear regression extrapolation (numpy.polyfit).\n\n"
            "4. 3-Period Moving Average:\n"
            "   Smooths short-term fluctuations to reveal underlying growth."
        ),
        title_color=C_ACCENT
    )

    # Right: Embedded chart
    if os.path.exists("forecast_chart.png"):
        slide5.shapes.add_picture("forecast_chart.png", Inches(6.4), Inches(1.6), width=Inches(6.1))

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 6: AGENTIC AI ARCHITECTURE (9 AGENTS)
    # ═════════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide6)
    add_header(slide6, "5. Agentic AI System: 9 Specialist Agents", "A team of autonomous AI agents collaborating to answer complex business questions.")

    rows, cols = 10, 3
    left, top, width, height = Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2)
    table_shape = slide6.shapes.add_table(rows, cols, left, top, width, height)
    t = table_shape.table
    t.columns[0].width = Inches(1.8)
    t.columns[1].width = Inches(4.2)
    t.columns[2].width = Inches(5.733)

    headers = ["Agent Name", "Primary Role & Focus", "Technical Mechanism / Output"]
    for j, h in enumerate(headers):
        cell = t.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_ACCENT
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

    agents_data = [
        ("🤖 Auto Router", "Intent Router & Pipeline Manager", "Analyzes user query and selects optimal agent sequence"),
        ("📊 Data Agent", "Database & SQL Expert", "Generates safe SELECT queries on Gold Star Schema tables"),
        ("🧠 ML Agent", "Statistician & Trend Analyst", "Calculates regression slopes & Z-score revenue anomalies"),
        ("🔮 Forecast Agent", "Predictor & Trend Extrapolator", "Extrapolates sales for future 1-12 months with Plotly lines"),
        ("💡 Insight Agent", "Business Analyst", "Transforms raw query data into 2-4 sentence evidence-backed insights"),
        ("🎯 Action Agent", "Management Consultant", "Recommends specific operational actions with business rationale"),
        ("📋 Report Agent", "Executive Assistant", "Summarizes full dashboard metrics into executive overviews"),
        ("⚙️ Pipeline Agent", "Data Engineer Audit Agent", "Queries metadata.pipeline_runs & dq_results for execution health"),
        ("💬 General Agent", "Technical Advisor", "Answers tech-stack questions on Airflow, Postgres, and Docker")
    ]

    for i, (name, role, mech) in enumerate(agents_data, start=1):
        for j, val in enumerate([name, role, mech]):
            cell = t.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if i % 2 == 1 else RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(9.5)
            p.font.color.rgb = C_PRIMARY

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 7: AGENT PIPELINE FLOW (WITH FLOWCHART DIAGRAM)
    # ═════════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide7)
    add_header(slide7, "6. Multi-Agent Orchestration Flow", "Autonomous intent classification, sequential agent chaining, and context forwarding.")

    # Embed agent flow diagram if available
    if os.path.exists("agent_flow_diagram.png"):
        slide7.shapes.add_picture("agent_flow_diagram.png", Inches(0.8), Inches(1.5), width=Inches(11.733))

    # Summary card below
    add_card(
        slide7, Inches(0.8), Inches(4.5), Inches(11.733), Inches(2.3),
        title="HOW AGENTS CHAIN TOGETHER",
        body=(
            "1. User Question → Auto Router analyzes intent and selects agent sequence (e.g. ['data', 'insight', 'action']).\n"
            "2. Data Agent executes safe SQL against Gold schema → Passes query results as context to Insight Agent.\n"
            "3. Insight Agent derives evidence-backed business insights → Passes insight as context to Action Agent.\n"
            "4. Action Agent formulates specific management recommendations → Streamlit UI renders formatted response cards."
        ),
        title_color=C_ACCENT
    )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 8: REAL QUESTION CASE STUDY
    # ═════════════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide8)
    add_header(slide8, "7. Case Study: Answering a Complex Question", "Step-by-step execution trace when a user asks a complex multi-part question.")

    add_card(
        slide8, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.8),
        title="",
        body="💬 User Question:  'Why did profit decrease and what should management focus on?'",
        title_color=C_PRIMARY
    )

    step_y = Inches(2.5)
    card_h = Inches(1.0)
    card_w = Inches(11.733)

    steps_info = [
        ("STEP 1: Auto Router 🤖", "Analyzes query intent → Selects agent execution chain: Data Agent 📊 → Insight Agent 💡 → Action Agent 🎯", C_ACCENT),
        ("STEP 2: Data Agent 📊", "Inspects Gold schema → Generates SELECT query → Executes safely → Retrieves revenue totals ($1.8M total sales).", C_SUCCESS),
        ("STEP 3: Insight Agent 💡", "Analyzes category numbers → Identifies: 'Category X sales declined by 12% MoM due to an 18% spike in regional freight costs.'", C_WARNING),
        ("STEP 4: Action Agent 🎯", "Formulates recommendation: 'Action: Renegotiate regional freight contracts and bundle Category X items with high-margin products.'", C_ACCENT),
    ]

    for idx, (title, desc, color) in enumerate(steps_info):
        add_card(
            slide8, Inches(0.8), step_y + Inches(idx * 1.15), card_w, card_h,
            title=title, body=desc, title_color=color
        )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 9: USER INTERFACE & DASHBOARD FEATURES
    # ═════════════════════════════════════════════════════════════════
    slide9 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide9)
    add_header(slide9, "8. Presentation Layer & UI Features", "Streamlit presentation layer with single-pane-of-glass executive controls.")

    add_card(
        slide9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        title="🧠 AI ANALYST CHATBOT TAB",
        body=(
            "• Interactive Chat Prompt: Ask questions in plain English.\n"
            "• 8 Suggested Prompt Chips: Clickable buttons ('Total sales?', 'Sales forecast?') for single-click answers.\n"
            "• Clear Chat Button: Reset conversation anytime.\n"
            "• Live Status Badges: Top-right green badges showing live status for Gemini (🟢), Database (🟢), and Pipeline (🟢).\n"
            "• Agent Pipeline Display: Shows active agent execution chain (e.g. Data → Insight → Action).\n"
            "• Visual Output: Generates Plotly charts, data tables, insight cards, and recommended action boxes."
        ),
        title_color=C_ACCENT
    )

    add_card(
        slide9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        title="📊 4 EXECUTIVE DASHBOARD PAGES",
        body=(
            "1. Overview Page:\n"
            "   6 KPI cards (Total Revenue, Orders, Customers, Products, AOV, Growth MoM) + Filter selects for Category, Region & Date.\n\n"
            "2. Pipeline Monitor:\n"
            "   Visual Medallion architecture flow + run logs from metadata.pipeline_runs.\n\n"
            "3. Data Quality Page:\n"
            "   Audit reports showing passed vs failed rules & failure distribution graphs.\n\n"
            "4. ML Insights Page:\n"
            "   Moving average lines, Z-score anomaly threshold graphs & MoM growth breakdown."
        ),
        title_color=C_SUCCESS
    )

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 10: CONCLUSION & EVALUATION DEMO
    # ═════════════════════════════════════════════════════════════════
    slide10 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide10)
    add_header(slide10, "9. Project Evaluation Summary & Live Demo", "100% operational, fully verified, and ready for live presentation.")

    add_card(
        slide10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2),
        title="⭐ KEY PROJECT HIGHLIGHTS",
        body=(
            "• End-to-End Pipeline: Raw data ingested, cleansed in Silver, and modeled into Gold Star Schema (3,938 records).\n"
            "• Historical Accuracy: Fully supports SCD Type 1 & Type 2 location history.\n"
            "• Robust Data Quality: Great Expectations rules & execution tracking.\n"
            "• Autonomous Multi-Agent AI: 9 specialized AI agents collaborating.\n"
            "• SQL Security: Enforces strict read-only query validator (blocks DROP, DELETE, etc.).\n"
            "• Rate-Limit Protection: Fast pattern fallbacks ensure instant answers (< 1 second)."
        ),
        title_color=C_SUCCESS
    )

    add_card(
        slide10, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2),
        title="🚀 LIVE EVALUATION DEMO URLS",
        body=(
            "1. Streamlit Application & AI Chatbot:\n"
            "   URL:  http://localhost:8501\n"
            "   Status:  100% Connected & Ready (All Badges Green 🟢)\n\n"
            "2. Apache Airflow Pipeline Orchestration:\n"
            "   URL:  http://localhost:8080\n"
            "   User: admin  /  Password: admin123\n"
            "   DAGs: 01_bronze_ingestion, 02_silver_transform, 03_gold_aggregation\n\n"
            "3. Comprehensive PDF Report & GitHub Repo:\n"
            "   GitHub: https://github.com/Eshaan4/E-COMMERCE-DATA-ENGINEERING-AGENTIC-AI-PLATFORM.git"
        ),
        title_color=C_ACCENT
    )

    prs.save(filename)
    print(f"Expanded 10-slide PowerPoint presentation successfully saved: {filename}")

if __name__ == "__main__":
    build_presentation()
